from __future__ import annotations

import io
import re
import time
from datetime import datetime, timezone
from typing import Any, Optional

from fastapi import Request

from open_webui.models.files import Files
from open_webui.models.skills import SkillModel
from open_webui.models.users import UserModel
from open_webui.storage.provider import Storage
from open_webui.utils.access_control import has_access
from open_webui.utils.server_files import save_server_file

try:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches, Pt
except ImportError:
    Presentation = None
    RGBColor = None
    MSO_SHAPE = None
    PP_ALIGN = None
    Inches = None
    Pt = None


BUILTIN_PPTX_SKILL_ID = "builtin:pptx-generator"
BUILTIN_PPTX_SKILL_IDENTIFIER = "halo.builtin.pptx-generator"
BUILTIN_PPTX_GENERATE_ENTRYPOINT_ID = "generate_pptx"
BUILTIN_PPTX_EDIT_ENTRYPOINT_ID = "edit_pptx"
BUILTIN_PPTX_ENTRYPOINT_ID = BUILTIN_PPTX_GENERATE_ENTRYPOINT_ID

PPTX_CONTENT_TYPE = (
    "application/vnd.openxmlformats-officedocument.presentationml.presentation"
)
MAX_SLIDES = 40
MAX_BULLETS_PER_SLIDE = 10
MAX_TEXT_CHARS = 24000
FONT_FAMILY = "Microsoft YaHei"


THEMES = {
    "blue": {
        "accent": (37, 99, 235),
        "accent_dark": (30, 64, 175),
        "background": (248, 250, 252),
        "surface": (255, 255, 255),
        "text": (15, 23, 42),
        "muted": (71, 85, 105),
    },
    "emerald": {
        "accent": (5, 150, 105),
        "accent_dark": (6, 95, 70),
        "background": (247, 254, 251),
        "surface": (255, 255, 255),
        "text": (15, 23, 42),
        "muted": (75, 85, 99),
    },
    "rose": {
        "accent": (225, 29, 72),
        "accent_dark": (159, 18, 57),
        "background": (255, 248, 250),
        "surface": (255, 255, 255),
        "text": (31, 41, 55),
        "muted": (107, 114, 128),
    },
    "slate": {
        "accent": (71, 85, 105),
        "accent_dark": (30, 41, 59),
        "background": (248, 250, 252),
        "surface": (255, 255, 255),
        "text": (15, 23, 42),
        "muted": (100, 116, 139),
    },
    "amber": {
        "accent": (217, 119, 6),
        "accent_dark": (146, 64, 14),
        "background": (255, 251, 235),
        "surface": (255, 255, 255),
        "text": (31, 41, 55),
        "muted": (92, 75, 56),
    },
}


def is_builtin_pptx_skill_id(value: Any) -> bool:
    return str(value or "").strip() in {
        BUILTIN_PPTX_SKILL_ID,
        BUILTIN_PPTX_SKILL_IDENTIFIER,
    }


def get_builtin_pptx_skill() -> SkillModel:
    timestamp = 0
    return SkillModel(
        id=BUILTIN_PPTX_SKILL_ID,
        user_id="system",
        name="PPTX Generator / Editor",
        description="服务端生成或编辑 PowerPoint .pptx 文件，并返回新的可下载文件。",
        content=(
            "Use this runnable skill when the user asks to create, draft, export, "
            "or edit a PowerPoint/PPTX presentation. For new decks, generate "
            "structured slide content and call generate_pptx. For an uploaded PPTX, "
            "use the file id from current_chat_resources as source_file_id and call "
            "edit_pptx. Always return the generated file to the user instead of "
            "printing raw JSON."
        ),
        source="builtin",
        identifier=BUILTIN_PPTX_SKILL_IDENTIFIER,
        source_url=None,
        meta={
            "kind": "skill_package",
            "builtin": True,
            "tags": ["pptx", "powerpoint", "presentation", "slides", "deck"],
            "manifest": {
                "name": "PPTX Generator / Editor",
                "identifier": BUILTIN_PPTX_SKILL_IDENTIFIER,
                "description": "Generate or edit downloadable PPTX files on the HaloWebUI server.",
                "category": "Documents",
                "tags": ["pptx", "slides", "presentation", "edit"],
            },
            "runtime": {
                "mode": "runnable",
                "entrypoints": [
                    {
                        "id": BUILTIN_PPTX_GENERATE_ENTRYPOINT_ID,
                        "runtime": "python",
                        "path": "builtin/pptx_generator.py",
                        "timeout": 60,
                        "description": (
                            "Generate a downloadable .pptx file on the server. "
                            "Use args like: {\"title\":\"Deck title\", "
                            "\"subtitle\":\"optional\", \"filename\":\"optional.pptx\", "
                            "\"theme\":\"blue|emerald|rose|slate|amber\", "
                            "\"slides\":[{\"title\":\"Slide title\", "
                            "\"bullets\":[\"point 1\",\"point 2\"], "
                            "\"body\":\"optional paragraph\", "
                            "\"notes\":\"optional speaker notes\"}]}. "
                            "If slides are not ready, pass markdown/outline/content text "
                            "and the skill will split it into slides. Returns JSON with "
                            "file.url and file.absolute_url."
                        ),
                    },
                    {
                        "id": BUILTIN_PPTX_EDIT_ENTRYPOINT_ID,
                        "runtime": "python",
                        "path": "builtin/pptx_editor.py",
                        "timeout": 60,
                        "description": (
                            "Edit an uploaded .pptx file and return a new downloadable "
                            ".pptx file. Required arg: source_file_id from "
                            "current_chat_resources. Use operations such as "
                            "[{\"type\":\"replace_text\",\"find\":\"old\",\"replace\":\"new\"}], "
                            "[{\"type\":\"add_slide\",\"title\":\"Summary\","
                            "\"bullets\":[\"point\"]}], "
                            "or [{\"type\":\"append_notes\",\"slide\":2,"
                            "\"notes\":\"speaker notes\"}]. "
                            "You may also pass replacements as a mapping and slides_to_add "
                            "or append_slides as slide objects. Do not overwrite the source file."
                        ),
                    },
                ],
                "install_status": "ready",
                "installed_hash": "builtin",
                "python_env_dir": None,
                "python_bin": None,
                "node_env_dir": None,
                "last_error": None,
                "installed_at": timestamp,
            },
        },
        access_control={},
        is_active=True,
        updated_at=timestamp,
        created_at=timestamp,
    )


def get_builtin_skill_by_id(skill_id: Any) -> Optional[SkillModel]:
    if is_builtin_pptx_skill_id(skill_id):
        return get_builtin_pptx_skill()
    return None


def list_builtin_chat_skills(user: Any = None) -> list[SkillModel]:
    return [get_builtin_pptx_skill()]


def _ensure_pptx_available() -> None:
    if (
        Presentation is None
        or RGBColor is None
        or MSO_SHAPE is None
        or PP_ALIGN is None
        or Inches is None
        or Pt is None
    ):
        raise RuntimeError("服务端缺少 python-pptx，无法生成 PPTX。")


def _rgb(value: tuple[int, int, int]):
    _ensure_pptx_available()
    return RGBColor(*value)


def _clean_text(value: Any, max_chars: int = 1200) -> str:
    text = str(value or "").replace("\r\n", "\n").replace("\r", "\n").strip()
    text = re.sub(r"[ \t]+", " ", text)
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text[:max_chars].strip()


def _clean_filename(value: Any, fallback_title: str) -> str:
    raw = _clean_text(value, 120) or fallback_title or "presentation"
    raw = re.sub(r"[<>:\"/\\|?*\x00-\x1f]", "_", raw)
    raw = re.sub(r"\s+", " ", raw).strip(" ._")
    if not raw:
        raw = "presentation"
    if not raw.lower().endswith(".pptx"):
        raw = f"{raw}.pptx"
    return raw[:140]


def _coerce_list(value: Any, *, max_items: int = MAX_BULLETS_PER_SLIDE) -> list[str]:
    if value is None:
        return []
    if isinstance(value, list):
        items = value
    elif isinstance(value, tuple):
        items = list(value)
    else:
        text = _clean_text(value, 3000)
        items = [
            re.sub(r"^\s*(?:[-*•]|\d+[\.)])\s+", "", line).strip()
            for line in text.splitlines()
            if line.strip()
        ]
    cleaned = [_clean_text(item, 420) for item in items]
    return [item for item in cleaned if item][:max_items]


def _new_slide(title: str = "", bullets: Optional[list[str]] = None) -> dict[str, Any]:
    return {"title": title, "bullets": bullets or [], "body": "", "notes": ""}


def _slide_from_mapping(value: dict[str, Any]) -> dict[str, Any]:
    title = _clean_text(
        value.get("title")
        or value.get("heading")
        or value.get("header")
        or value.get("name"),
        180,
    )
    body = _clean_text(
        value.get("body") or value.get("content") or value.get("text") or "",
        1500,
    )
    bullets = _coerce_list(
        value.get("bullets")
        or value.get("points")
        or value.get("items")
        or value.get("key_points")
    )
    if not bullets and body:
        body_lines = _coerce_list(body)
        if len(body_lines) > 1:
            bullets = body_lines
            body = ""
    return {
        "title": title or "Untitled",
        "subtitle": _clean_text(value.get("subtitle") or "", 220),
        "body": body,
        "bullets": bullets,
        "notes": _clean_text(value.get("notes") or value.get("speaker_notes") or "", 1200),
    }


def _parse_markdown_slides(content: str) -> list[dict[str, Any]]:
    content = _clean_text(content, MAX_TEXT_CHARS)
    if not content:
        return []

    slides: list[dict[str, Any]] = []
    current: Optional[dict[str, Any]] = None

    def flush_current() -> None:
        nonlocal current
        if current and (current.get("title") or current.get("bullets") or current.get("body")):
            slides.append(current)
        current = None

    for raw_line in content.splitlines():
        line = raw_line.strip()
        if not line:
            continue

        heading = re.match(r"^(#{1,3})\s+(.+)$", line)
        if heading:
            flush_current()
            current = _new_slide(_clean_text(heading.group(2), 180))
            continue

        bullet = re.match(r"^(?:[-*•]|\d+[\.)])\s+(.+)$", line)
        if current is None:
            current = _new_slide("Overview")
        if bullet:
            current.setdefault("bullets", []).append(_clean_text(bullet.group(1), 420))
        else:
            current.setdefault("bullets", []).append(_clean_text(line, 420))

    flush_current()
    return _rebalance_slides(slides)


def _rebalance_slides(slides: list[dict[str, Any]]) -> list[dict[str, Any]]:
    balanced: list[dict[str, Any]] = []
    for slide in slides:
        bullets = _coerce_list(slide.get("bullets"), max_items=100)
        if len(bullets) <= MAX_BULLETS_PER_SLIDE:
            slide["bullets"] = bullets
            balanced.append(slide)
            continue

        title = _clean_text(slide.get("title"), 180) or "Overview"
        for index in range(0, len(bullets), MAX_BULLETS_PER_SLIDE):
            suffix = "" if index == 0 else f" ({index // MAX_BULLETS_PER_SLIDE + 1})"
            balanced.append(
                {
                    **slide,
                    "title": f"{title}{suffix}",
                    "bullets": bullets[index : index + MAX_BULLETS_PER_SLIDE],
                    "body": "" if index else slide.get("body", ""),
                }
            )
    return balanced[:MAX_SLIDES]


def _normalize_slides(args: dict[str, Any]) -> list[dict[str, Any]]:
    raw_slides = args.get("slides")
    slides: list[dict[str, Any]] = []
    if isinstance(raw_slides, list):
        for item in raw_slides:
            if isinstance(item, dict):
                slides.append(_slide_from_mapping(item))
            elif str(item or "").strip():
                slides.append(_new_slide(_clean_text(item, 180)))

    if not slides:
        outline = (
            args.get("outline")
            or args.get("markdown")
            or args.get("content")
            or args.get("text")
            or ""
        )
        slides = _parse_markdown_slides(str(outline or ""))

    if not slides:
        slides = [
            _new_slide(
                "Overview",
                [
                    "Clarify the target audience and goal.",
                    "Summarize the core message.",
                    "Add supporting evidence and next steps.",
                ],
            )
        ]

    return _rebalance_slides(slides)[:MAX_SLIDES]


def _theme_from_args(args: dict[str, Any]) -> dict[str, Any]:
    raw_theme = args.get("theme")
    theme_name = ""
    if isinstance(raw_theme, dict):
        theme_name = str(raw_theme.get("name") or raw_theme.get("preset") or "").lower()
    else:
        theme_name = str(raw_theme or "").lower()
    theme = THEMES.get(theme_name, THEMES["blue"])
    return {key: _rgb(value) for key, value in theme.items()}


def _set_background(slide, color: Any) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _set_run(run, *, size: float, color: Any, bold: bool = False) -> None:
    run.font.name = FONT_FAMILY
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def _add_text_box(
    slide,
    *,
    left,
    top,
    width,
    height,
    text: str,
    size: float,
    color: Any,
    bold: bool = False,
    align: Any = None,
) -> None:
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align if align is not None else PP_ALIGN.LEFT
    run = paragraph.add_run()
    run.text = text
    _set_run(run, size=size, color=color, bold=bold)


def _add_cover(prs: Any, title: str, subtitle: str, theme: dict[str, Any]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_background(slide, theme["background"])
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.18)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = theme["accent"]
    bar.line.color.rgb = theme["accent"]
    rail = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.65), Inches(1.15), Inches(0.12), Inches(4.75)
    )
    rail.fill.solid()
    rail.fill.fore_color.rgb = theme["accent"]
    rail.line.color.rgb = theme["accent"]

    _add_text_box(
        slide,
        left=Inches(1.05),
        top=Inches(1.5),
        width=Inches(10.8),
        height=Inches(1.4),
        text=title,
        size=40,
        color=theme["text"],
        bold=True,
    )
    if subtitle:
        _add_text_box(
            slide,
            left=Inches(1.08),
            top=Inches(3.1),
            width=Inches(9.8),
            height=Inches(0.85),
            text=subtitle,
            size=20,
            color=theme["muted"],
        )
    _add_text_box(
        slide,
        left=Inches(1.08),
        top=Inches(6.55),
        width=Inches(6.0),
        height=Inches(0.35),
        text=datetime.now(timezone.utc).strftime("%Y-%m-%d"),
        size=10,
        color=theme["muted"],
    )


def _add_agenda(prs: Any, slides: list[dict[str, Any]], theme: dict[str, Any]) -> None:
    agenda_items = [_clean_text(slide.get("title"), 120) for slide in slides[:8]]
    agenda_items = [item for item in agenda_items if item]
    if len(agenda_items) < 3:
        return

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_background(slide, theme["background"])
    _add_text_box(
        slide,
        left=Inches(0.75),
        top=Inches(0.55),
        width=Inches(11.4),
        height=Inches(0.7),
        text="Agenda",
        size=28,
        color=theme["text"],
        bold=True,
    )
    _add_bullets(
        slide,
        bullets=agenda_items,
        left=Inches(1.0),
        top=Inches(1.65),
        width=Inches(10.8),
        height=Inches(4.8),
        theme=theme,
        font_size=22,
    )


def _add_bullets(
    slide,
    *,
    bullets: list[str],
    left,
    top,
    width,
    height,
    theme: dict[str, Any],
    font_size: float = 19,
) -> None:
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True

    for index, item in enumerate(bullets[:MAX_BULLETS_PER_SLIDE]):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.space_after = Pt(8)
        run = paragraph.add_run()
        run.text = f"• {item}"
        _set_run(run, size=font_size, color=theme["text"])


def _add_content_slide(
    prs: Any,
    slide_data: dict[str, Any],
    slide_number: int,
    total_slides: int,
    theme: dict[str, Any],
) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_background(slide, theme["background"])

    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.11)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = theme["accent"]
    bar.line.color.rgb = theme["accent"]

    title = _clean_text(slide_data.get("title"), 180) or f"Slide {slide_number}"
    _add_text_box(
        slide,
        left=Inches(0.75),
        top=Inches(0.45),
        width=Inches(11.2),
        height=Inches(0.65),
        text=title,
        size=25,
        color=theme["text"],
        bold=True,
    )

    subtitle = _clean_text(slide_data.get("subtitle"), 220)
    top = 1.35
    if subtitle:
        _add_text_box(
            slide,
            left=Inches(0.78),
            top=Inches(1.08),
            width=Inches(10.8),
            height=Inches(0.4),
            text=subtitle,
            size=13,
            color=theme["muted"],
        )
        top = 1.65

    bullets = _coerce_list(slide_data.get("bullets"))
    body = _clean_text(slide_data.get("body"), 1500)
    if bullets:
        _add_bullets(
            slide,
            bullets=bullets,
            left=Inches(0.95),
            top=Inches(top),
            width=Inches(11.1),
            height=Inches(4.8),
            theme=theme,
        )
    elif body:
        _add_text_box(
            slide,
            left=Inches(0.95),
            top=Inches(top),
            width=Inches(10.8),
            height=Inches(4.8),
            text=body,
            size=18,
            color=theme["text"],
        )

    _add_text_box(
        slide,
        left=Inches(11.7),
        top=Inches(6.78),
        width=Inches(0.9),
        height=Inches(0.28),
        text=f"{slide_number}/{total_slides}",
        size=9,
        color=theme["muted"],
        align=PP_ALIGN.RIGHT,
    )

    notes = _clean_text(slide_data.get("notes"), 1200)
    if notes:
        slide.notes_slide.notes_text_frame.text = notes


def _iter_text_frames(slide: Any):
    for shape in getattr(slide, "shapes", []):
        if getattr(shape, "has_text_frame", False):
            yield shape.text_frame
        elif getattr(shape, "has_table", False):
            for row in shape.table.rows:
                for cell in row.cells:
                    yield cell.text_frame


def _replace_text_in_prs(
    prs: Any, find: str, replace: str, slide_indexes: list[int]
) -> int:
    search = _clean_text(find, 240)
    if not search:
        return 0

    replacement = str(replace or "")
    replaced = 0
    for index, slide in enumerate(prs.slides):
        if slide_indexes and index not in slide_indexes:
            continue
        for text_frame in _iter_text_frames(slide):
            for paragraph in text_frame.paragraphs:
                for run in paragraph.runs:
                    if search in run.text:
                        run.text = run.text.replace(search, replacement)
                        replaced += 1
    return replaced


def _append_notes_to_slide(slide: Any, notes: str) -> bool:
    cleaned_notes = _clean_text(notes, 2000)
    if not cleaned_notes:
        return False

    notes_frame = slide.notes_slide.notes_text_frame
    existing = _clean_text(getattr(notes_frame, "text", "") or "", 2000)
    notes_frame.text = (
        f"{existing}\n{cleaned_notes}".strip() if existing else cleaned_notes
    )
    return True


def _coerce_slide_indexes(value: Any, total_slides: int) -> list[int]:
    if value is None:
        return list(range(total_slides))

    values = value if isinstance(value, (list, tuple, set)) else [value]
    indexes: list[int] = []
    for item in values:
        try:
            index = int(item)
        except Exception:
            continue
        if index <= 0 or index > total_slides:
            continue
        zero_based = index - 1
        if zero_based not in indexes:
            indexes.append(zero_based)
    return indexes


def _normalize_edit_operations(args: dict[str, Any]) -> list[dict[str, Any]]:
    operations: list[dict[str, Any]] = []

    raw_operations = args.get("operations")
    if isinstance(raw_operations, list):
        for item in raw_operations:
            if isinstance(item, dict):
                operations.append(item)

    replacements = args.get("replacements")
    if isinstance(replacements, dict):
        for find, replace in replacements.items():
            operations.append(
                {
                    "type": "replace_text",
                    "find": find,
                    "replace": replace,
                }
            )
    elif isinstance(replacements, list):
        for item in replacements:
            if not isinstance(item, dict):
                continue
            operations.append(
                {
                    "type": "replace_text",
                    "find": item.get("find") or item.get("old"),
                    "replace": item.get("replace") or item.get("new") or "",
                    **(
                        {"slide": item.get("slide")}
                        if item.get("slide") is not None
                        else {}
                    ),
                }
            )

    for key in ("slides_to_add", "append_slides"):
        extra_slides = args.get(key)
        if isinstance(extra_slides, list):
            for item in extra_slides:
                if isinstance(item, dict):
                    operations.append({"type": "add_slide", **item})

    if isinstance(args.get("slides"), list) and not operations:
        for item in args.get("slides") or []:
            if isinstance(item, dict):
                operations.append({"type": "add_slide", **item})

    return operations


def _apply_edit_operations(prs: Any, args: dict[str, Any]) -> dict[str, Any]:
    operations = _normalize_edit_operations(args)
    if not operations:
        raise RuntimeError("未提供可执行的编辑操作。")

    applied = 0
    replaced_total = 0
    notes_total = 0
    added_slides = 0
    theme = (
        _theme_from_args(args)
        if args.get("theme")
        else {key: _rgb(value) for key, value in THEMES["slate"].items()}
    )

    for operation in operations:
        op_type = _clean_text(operation.get("type"), 80).lower().replace("-", "_")
        if op_type in {"replace_text", "replace"}:
            find = (
                operation.get("find")
                or operation.get("old")
                or operation.get("search")
            )
            replace = operation.get("replace") or operation.get("new") or ""
            slide_target = operation.get("slide")
            slide_indexes = _coerce_slide_indexes(slide_target, len(prs.slides))
            if slide_target is not None and not slide_indexes:
                continue
            replaced = _replace_text_in_prs(
                prs, str(find or ""), str(replace or ""), slide_indexes
            )
            if replaced > 0:
                applied += 1
                replaced_total += replaced
            continue

        if op_type in {"append_notes", "add_notes"}:
            slide_target = operation.get("slide")
            slide_indexes = _coerce_slide_indexes(slide_target, len(prs.slides))
            if slide_target is not None and not slide_indexes:
                continue
            note_text = str(operation.get("notes") or operation.get("text") or "")
            for slide_index in slide_indexes or []:
                if 0 <= slide_index < len(prs.slides):
                    if _append_notes_to_slide(prs.slides[slide_index], note_text):
                        notes_total += 1
                        applied += 1
            continue

        if op_type in {"add_slide", "append_slide", "insert_slide"}:
            slide_data = (
                operation.get("slide")
                if isinstance(operation.get("slide"), dict)
                else operation
            )
            if not isinstance(slide_data, dict):
                slide_data = {}
            slide_number = len(prs.slides) + 1
            _add_content_slide(
                prs,
                _slide_from_mapping(slide_data),
                slide_number,
                slide_number,
                theme,
            )
            added_slides += 1
            applied += 1
            continue

        if op_type in {"add_slides", "append_slides"}:
            slides = operation.get("slides")
            if isinstance(slides, list):
                for slide_data in slides:
                    if not isinstance(slide_data, dict):
                        continue
                    slide_number = len(prs.slides) + 1
                    _add_content_slide(
                        prs,
                        _slide_from_mapping(slide_data),
                        slide_number,
                        slide_number,
                        theme,
                    )
                    added_slides += 1
                    applied += 1
            continue

    return {
        "applied_operations": applied,
        "replaced_count": replaced_total,
        "notes_count": notes_total,
        "added_slides": added_slides,
    }


def _save_pptx_file(
    request: Request,
    user: UserModel,
    pptx_bytes: bytes,
    filename: str,
    *,
    skill_source: str,
    source_file_id: Optional[str] = None,
    edit_stats: Optional[dict[str, Any]] = None,
) -> tuple[str, dict[str, Any]]:
    attachment = save_server_file(
        request,
        user,
        pptx_bytes,
        filename,
        PPTX_CONTENT_TYPE,
        producer=skill_source,
        metadata={
            "skill_id": BUILTIN_PPTX_SKILL_ID,
            "created_from_skill_id": BUILTIN_PPTX_SKILL_ID,
            **(
                {"source_file_id": source_file_id}
                if source_file_id
                else {}
            ),
            **({"edit_stats": edit_stats} if edit_stats else {}),
        },
        preview={"kind": "pptx", "strategy": "client_ooxml"},
    )

    return attachment["id"], attachment


def _load_source_pptx(
    request: Request, user: UserModel, source_file_id: Any
) -> tuple[bytes, str]:
    file_id = _clean_text(source_file_id, 120)
    if not file_id:
        raise RuntimeError("未提供 source_file_id。")

    file_item = Files.get_file_by_id(file_id)
    if not file_item:
        raise RuntimeError("找不到要编辑的源 PPTX 文件。")
    if (
        file_item.user_id != user.id
        and getattr(user, "role", "") != "admin"
        and not has_access(user.id, "read", file_item.access_control)
    ):
        raise RuntimeError("当前用户没有权限编辑这个文件。")

    file_name = _clean_text(getattr(file_item, "filename", "") or "", 240)
    content_type = str((file_item.meta or {}).get("content_type") or "").lower()
    if not file_name.lower().endswith(".pptx") and content_type not in {
        PPTX_CONTENT_TYPE,
        "application/vnd.ms-powerpoint",
    }:
        raise RuntimeError("源文件不是可编辑的 PPTX。")

    if not getattr(file_item, "path", None):
        raise RuntimeError("源文件缺少存储路径。")

    file_path = Storage.get_file(file_item.path)
    with open(file_path, "rb") as file_handle:
        return file_handle.read(), file_name


def _presentation_bytes(prs: Any) -> bytes:
    buffer = io.BytesIO()
    prs.save(buffer)
    return buffer.getvalue()


def generate_pptx_bytes(args: dict[str, Any]) -> tuple[bytes, str, int]:
    args = args if isinstance(args, dict) else {}
    slides = _normalize_slides(args)
    title = _clean_text(args.get("title"), 160) or _clean_text(
        slides[0].get("title"), 160
    )
    title = title or "Presentation"
    subtitle = _clean_text(args.get("subtitle") or args.get("description"), 260)
    filename = _clean_filename(args.get("filename"), title)
    theme = _theme_from_args(args)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    _add_cover(prs, title, subtitle, theme)
    if bool(args.get("include_agenda", len(slides) >= 3)):
        _add_agenda(prs, slides, theme)

    total_content_slides = len(slides)
    for index, slide in enumerate(slides, start=1):
        _add_content_slide(prs, slide, index, total_content_slides, theme)

    buffer = io.BytesIO()
    prs.save(buffer)
    return buffer.getvalue(), filename, len(prs.slides)


def edit_pptx_bytes(
    source_pptx_bytes: bytes,
    args: dict[str, Any],
    source_filename: str = "presentation.pptx",
) -> tuple[bytes, str, int, dict[str, Any]]:
    _ensure_pptx_available()
    args = args if isinstance(args, dict) else {}
    prs = Presentation(io.BytesIO(source_pptx_bytes))
    edit_stats = _apply_edit_operations(prs, args)
    filename = _clean_filename(
        args.get("filename"), _clean_text(source_filename, 160) or "presentation"
    )
    return _presentation_bytes(prs), filename, len(prs.slides), edit_stats


def create_pptx_file(
    request: Request,
    user: UserModel,
    args: dict[str, Any],
) -> dict[str, Any]:
    try:
        pptx_bytes, filename, slide_count = generate_pptx_bytes(args)
    except ImportError as exc:
        raise RuntimeError("服务端缺少 python-pptx，无法生成 PPTX。") from exc

    now = int(time.time())
    _file_id, file_result = _save_pptx_file(
        request,
        user,
        pptx_bytes,
        filename,
        skill_source="pptx_generator",
    )

    return {
        "ok": True,
        "skill_id": BUILTIN_PPTX_SKILL_ID,
        "entrypoint_id": BUILTIN_PPTX_GENERATE_ENTRYPOINT_ID,
        "created_at": now,
        "slide_count": slide_count,
        "file": file_result,
        "files": [file_result],
        "message": "PPTX generated and attached to the chat message.",
    }


def create_pptx_edit_file(
    request: Request,
    user: UserModel,
    args: dict[str, Any],
) -> dict[str, Any]:
    args = args if isinstance(args, dict) else {}
    source_file_id = (
        args.get("source_file_id")
        or args.get("file_id")
        or args.get("pptx_file_id")
        or args.get("input_file_id")
    )
    source_bytes, source_filename = _load_source_pptx(request, user, source_file_id)
    edited_bytes, filename, slide_count, edit_stats = edit_pptx_bytes(
        source_bytes,
        args,
        source_filename=source_filename,
    )

    now = int(time.time())
    _, file_result = _save_pptx_file(
        request,
        user,
        edited_bytes,
        filename,
        skill_source="pptx_editor",
        source_file_id=_clean_text(source_file_id, 120),
        edit_stats=edit_stats,
    )

    return {
        "ok": True,
        "skill_id": BUILTIN_PPTX_SKILL_ID,
        "entrypoint_id": BUILTIN_PPTX_EDIT_ENTRYPOINT_ID,
        "created_at": now,
        "slide_count": slide_count,
        "edit_stats": edit_stats,
        "source_file_id": _clean_text(source_file_id, 120),
        "file": file_result,
        "files": [file_result],
        "message": "PPTX edited and attached to the chat message.",
    }
